import os
import tempfile
import unittest
from datetime import datetime, timedelta, timezone
from pathlib import Path
from zipfile import ZipFile

os.environ["DATABASE_URL"] = f"sqlite:///{tempfile.NamedTemporaryFile(delete=False).name}"

import api.documents as documents_api
from app.features.documents.generation import create_generated_file
from app.features.documents.renderer import render_docx, render_document
from app.features.documents.workspace_save import save_generated_file_to_workspace
from app.features.workspaces.files.storage import WorkspaceStorageConfig
from fastapi import HTTPException
from fastapi.responses import FileResponse
from models import Base, SessionLocal, engine
from models.generated_file import GeneratedFile
from models.user import User
from models.workspace import Workspace, WorkspaceFile


class DocumentsPhase11Tests(unittest.TestCase):
    def setUp(self):
        Base.metadata.drop_all(bind=engine)
        Base.metadata.create_all(bind=engine)
        self.db = SessionLocal()
        self.temp = tempfile.TemporaryDirectory()
        self.user = User(username="doc", password_hash="hash", role="employee", nickname="Doc")
        self.admin = User(username="admin-doc", password_hash="hash", role="admin", nickname="Admin")
        self.other = User(username="other-doc", password_hash="hash", role="employee", nickname="Other")
        self.db.add_all([self.user, self.admin, self.other])
        self.db.commit()
        self.db.refresh(self.user)
        self.db.refresh(self.admin)
        self.db.refresh(self.other)

    def tearDown(self):
        self.temp.cleanup()
        self.db.close()

    def create_file_record(self, user_id: int, file_id: str = "file-1", age_hours: int = 1):
        path = Path(self.temp.name) / f"{file_id}.docx"
        path.write_bytes(b"docx")
        record = GeneratedFile(
            id=file_id,
            user_id=user_id,
            session_id=None,
            filename=f"{file_id}.docx",
            path=str(path),
            mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            created_at=datetime.now(timezone.utc) - timedelta(hours=age_hours),
        )
        self.db.add(record)
        self.db.commit()
        return record

    def test_download_generated_file_requires_owner(self):
        record = self.create_file_record(self.user.id)

        response = documents_api.download_generated_file(record.id, self.user, self.db)

        self.assertIsInstance(response, FileResponse)
        with self.assertRaises(HTTPException) as exc:
            documents_api.download_generated_file(record.id, self.other, self.db)
        self.assertEqual(exc.exception.status_code, 404)

    def test_cleanup_generated_files_is_admin_only_and_removes_expired(self):
        expired = self.create_file_record(self.user.id, "expired", age_hours=72)
        fresh = self.create_file_record(self.user.id, "fresh", age_hours=1)

        with self.assertRaises(HTTPException) as exc:
            documents_api.cleanup_generated_files(self.user, self.db)

        self.assertEqual(exc.exception.status_code, 403)
        result = documents_api.cleanup_generated_files(self.admin, self.db)

        self.assertEqual(result["removed"], 1)
        self.assertIsNone(self.db.get(GeneratedFile, expired.id))
        self.assertIsNotNone(self.db.get(GeneratedFile, fresh.id))
        self.assertFalse(Path(expired.path).exists())
        self.assertTrue(Path(fresh.path).exists())

    def test_docx_renderer_strips_markdown_wrappers_and_source_list(self):
        path = Path(self.temp.name) / "rendered.docx"
        render_docx(
            "用车申请",
            (
                "说明：下面是模板。\n\n"
                "```markdown\n"
                "# 用车申请\n\n"
                "**申请人**：张三\n\n"
                "| 字段 | 内容 |\n"
                "| --- | --- |\n"
                "| 用车类型 | 临时用车 |\n"
                "\n"
                "```\n\n"
                "本次回答使用的来源文件：\n"
                "- [[用车申请]]\n"
            ),
            path,
        )

        with ZipFile(path) as archive:
            document_xml = archive.read("word/document.xml").decode("utf-8")

        self.assertIn("用车申请", document_xml)
        self.assertIn("申请人", document_xml)
        self.assertNotIn("```", document_xml)
        self.assertNotIn("**", document_xml)
        self.assertNotIn("[[用车申请]]", document_xml)

    def test_renderer_registry_supports_markdown_and_plain_text(self):
        md_path = Path(self.temp.name) / "result.md"
        txt_path = Path(self.temp.name) / "result.txt"

        render_document("markdown", "会议纪要", "## 决策\n\n**继续推进**", md_path)
        render_document("txt", "会议纪要", "## 决策\n\n**继续推进**", txt_path)

        self.assertIn("# 会议纪要", md_path.read_text(encoding="utf-8"))
        self.assertIn("## 决策", md_path.read_text(encoding="utf-8"))
        text = txt_path.read_text(encoding="utf-8")
        self.assertIn("决策", text)
        self.assertIn("继续推进", text)
        self.assertNotIn("**", text)

        with self.assertRaises(ValueError):
            render_document("eml", "x", "y", Path(self.temp.name) / "bad.eml")

    def test_renderer_registry_supports_xlsx_and_pptx(self):
        from openpyxl import load_workbook
        from pptx import Presentation

        xlsx_path = Path(self.temp.name) / "result.xlsx"
        pptx_path = Path(self.temp.name) / "result.pptx"
        content = (
            "# 项目报价\n\n"
            "| 项目 | 金额 |\n"
            "| --- | --- |\n"
            "| 门窗 | 12000 |\n"
            "| 幕墙 | 23000 |\n"
            "\n"
            "## 汇报要点\n"
            "- 报价结构已整理\n"
            "- 风险项需要复核\n"
        )

        render_document("xlsx", "项目报价", content, xlsx_path)
        render_document("pptx", "项目报价", content, pptx_path)

        workbook = load_workbook(xlsx_path)
        worksheet = workbook.active
        self.assertEqual(worksheet.cell(row=1, column=1).value, "项目")
        self.assertEqual(worksheet.cell(row=2, column=1).value, "门窗")
        self.assertEqual(worksheet.cell(row=3, column=2).value, 23000)

        presentation = Presentation(pptx_path)
        self.assertGreaterEqual(len(presentation.slides), 1)
        slide_text = "\n".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        self.assertIn("项目报价", slide_text)
        self.assertIn("报价结构已整理", slide_text)

    def test_xlsx_renderer_splits_multiple_markdown_tables_into_sheets(self):
        from openpyxl import load_workbook

        path = Path(self.temp.name) / "multi.xlsx"
        render_document(
            "xlsx",
            "项目数据",
            (
                "| 名称 | 数量 |\n"
                "| --- | --- |\n"
                "| 门 | 12 |\n\n"
                "| 风险 | 等级 |\n"
                "| --- | --- |\n"
                "| 交期 | 3 |\n"
            ),
            path,
        )

        workbook = load_workbook(path)
        self.assertEqual(len(workbook.sheetnames), 2)
        self.assertEqual(workbook[workbook.sheetnames[0]].freeze_panes, "A2")
        self.assertEqual(workbook[workbook.sheetnames[0]].cell(row=2, column=2).value, 12)
        self.assertEqual(workbook[workbook.sheetnames[1]].cell(row=2, column=1).value, "交期")

    def test_pptx_renderer_splits_long_sections_into_continuation_slides(self):
        from pptx import Presentation

        path = Path(self.temp.name) / "long.pptx"
        render_document(
            "pptx",
            "项目周报",
            "# 风险清单\n\n" + "\n".join(f"- 风险项 {index}" for index in range(1, 9)),
            path,
        )

        presentation = Presentation(path)
        slide_text = "\n".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        self.assertGreaterEqual(len(presentation.slides), 3)
        self.assertIn("风险清单（续）", slide_text)
        self.assertNotIn("**", slide_text)

    def test_pdf_renderer_supports_text_formats_without_markdown_noise(self):
        from pypdf import PdfReader

        path = Path(self.temp.name) / "result.pdf"
        render_document(
            "pdf",
            "会议纪要",
            "# 会议纪要\n\n## 决策\n\n**继续推进**\n\n| 字段 | 内容 |\n| --- | --- |\n| 负责人 | 张三 |",
            path,
        )

        reader = PdfReader(str(path))
        self.assertGreaterEqual(len(reader.pages), 1)
        extracted = "\n".join(page.extract_text() or "" for page in reader.pages)
        self.assertIn("会议纪要", extracted)
        self.assertIn("继续推进", extracted)
        self.assertNotIn("**", extracted)
        self.assertNotIn("| --- |", extracted)

    def test_create_generated_file_uses_format_metadata(self):
        payload = create_generated_file(
            self.db,
            self.user.id,
            None,
            "项目总结",
            "# 项目总结\n\n内容",
            output_format="md",
            generated_files_root=Path(self.temp.name),
        )

        record = self.db.get(GeneratedFile, payload["id"])
        self.assertIsNotNone(record)
        self.assertTrue(record.filename.endswith(".md"))
        self.assertEqual(record.mime_type, "text/markdown; charset=utf-8")
        self.assertTrue(Path(record.path).exists())

    def test_create_generated_file_uses_office_format_metadata(self):
        payload = create_generated_file(
            self.db,
            self.user.id,
            None,
            "项目汇报",
            "- 第一页\n- 第二页",
            output_format="ppt",
            generated_files_root=Path(self.temp.name),
        )

        record = self.db.get(GeneratedFile, payload["id"])
        self.assertIsNotNone(record)
        self.assertTrue(record.filename.endswith(".pptx"))
        self.assertEqual(record.mime_type, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        self.assertTrue(Path(record.path).exists())

    def test_create_generated_file_uses_pdf_format_metadata(self):
        payload = create_generated_file(
            self.db,
            self.user.id,
            None,
            "会议纪要",
            "# 会议纪要\n\n- 决策",
            output_format="pdf",
            generated_files_root=Path(self.temp.name),
        )

        record = self.db.get(GeneratedFile, payload["id"])
        self.assertIsNotNone(record)
        self.assertTrue(record.filename.endswith(".pdf"))
        self.assertEqual(record.mime_type, "application/pdf")
        self.assertTrue(Path(record.path).exists())

    def test_save_generated_file_to_project_workspace_unfiled_dir(self):
        workspace_root = Path(self.temp.name) / "project" / "BFI" / "test-project"
        workspace = Workspace(
            name="Test Project",
            slug="test-project",
            created_by=self.user.id,
            storage_path=str(workspace_root),
            brand="BFI",
            workspace_kind="project",
        )
        self.db.add(workspace)
        self.db.commit()
        self.db.refresh(workspace)
        payload = create_generated_file(
            self.db,
            self.user.id,
            None,
            "报价草稿",
            "内容",
            output_format="txt",
            generated_files_root=Path(self.temp.name) / "generated",
        )

        result = save_generated_file_to_workspace(
            self.db,
            workspace=workspace,
            user=self.user,
            generated_file_id=payload["id"],
            conflict_strategy="keep_both",
            storage_config=WorkspaceStorageConfig(
                workspaces_root=Path(self.temp.name),
                project_root_name="project",
                customer_root_name="customer",
                project_brands=("BFI",),
                customer_brand="CUSTOMER",
                crm_workspace_slug="CRM",
                crm_raw_dir="raw",
            ),
        )
        self.db.commit()

        self.assertTrue(result["ok"])
        self.assertEqual(result["path"], "99-未归档文件/报价草稿.txt")
        self.assertTrue((workspace_root / result["path"]).exists())
        meta = self.db.get(WorkspaceFile, result["file_id"])
        self.assertEqual(meta.content_type, "text/plain; charset=utf-8")

    def test_save_generated_file_rejects_user_workspace(self):
        workspace = Workspace(
            name="User Space",
            slug="user-space",
            created_by=self.user.id,
            storage_path=str(Path(self.temp.name) / "user"),
            brand="BFI",
            workspace_kind="user",
        )
        self.db.add(workspace)
        self.db.commit()
        payload = create_generated_file(
            self.db,
            self.user.id,
            None,
            "个人草稿",
            "内容",
            output_format="txt",
            generated_files_root=Path(self.temp.name) / "generated",
        )

        with self.assertRaises(HTTPException) as exc:
            save_generated_file_to_workspace(
                self.db,
                workspace=workspace,
                user=self.user,
                generated_file_id=payload["id"],
                conflict_strategy="keep_both",
                storage_config=WorkspaceStorageConfig(
                    workspaces_root=Path(self.temp.name),
                    project_root_name="project",
                    customer_root_name="customer",
                    project_brands=("BFI",),
                    customer_brand="CUSTOMER",
                    crm_workspace_slug="CRM",
                    crm_raw_dir="raw",
                ),
            )

        self.assertEqual(exc.exception.status_code, 400)


if __name__ == "__main__":
    unittest.main()
