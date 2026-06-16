from __future__ import annotations

from pathlib import Path

from app.features.documents.content_parser import DocumentBlock, ParsedDocument, parse_document


MAX_ITEMS_PER_SLIDE = 6


def render_pptx(title: str, content: str, output_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Pt

    parsed = parse_document(title, content)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.core_properties.title = parsed.title

    _add_slide(presentation, parsed.title, _first_slide_items(parsed), is_title=True)
    for slide_title, items in _outline_slides(parsed):
        for chunk_index in range(0, len(items), MAX_ITEMS_PER_SLIDE):
            chunk = items[chunk_index:chunk_index + MAX_ITEMS_PER_SLIDE]
            title_suffix = "" if chunk_index == 0 else "（续）"
            _add_slide(presentation, f"{slide_title}{title_suffix}", chunk)

    presentation.save(output_path)
    return output_path


def _first_slide_items(parsed: ParsedDocument) -> list[str]:
    items = [block.text for block in parsed.blocks if block.type in {"paragraph", "bullet", "numbered"} and block.text]
    return items[:3] or [parsed.title]


def _outline_slides(parsed: ParsedDocument) -> list[tuple[str, list[str]]]:
    slides: list[tuple[str, list[str]]] = []
    current_title = parsed.title
    current_items: list[str] = []
    for block in parsed.blocks:
        if block.type == "heading":
            if current_items:
                slides.append((current_title, current_items))
                current_items = []
            current_title = block.text
            continue
        if block.type == "table":
            current_items.extend(_table_to_items(block))
        elif block.type in {"paragraph", "bullet", "numbered", "quote"} and block.text:
            current_items.append(block.text)
    if current_items:
        slides.append((current_title, current_items))
    return slides[:12]


def _table_to_items(block: DocumentBlock) -> list[str]:
    if not block.rows:
        return []
    header = block.rows[0]
    items = []
    for row in block.rows[1:]:
        pairs = [f"{header[index]}: {value}" for index, value in enumerate(row) if index < len(header)]
        items.append("；".join(pairs))
    return items


def _add_slide(presentation, title: str, items: list[str], *, is_title: bool = False) -> None:
    from pptx.util import Pt

    layout = presentation.slide_layouts[0] if is_title else presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(layout)
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = title
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(34 if is_title else 28)
    placeholders = [shape for shape in slide.placeholders if shape != title_shape]
    if placeholders:
        body = placeholders[0].text_frame
        body.clear()
        for item_index, item in enumerate(items):
            paragraph = body.paragraphs[0] if item_index == 0 else body.add_paragraph()
            paragraph.text = item
            paragraph.level = 0
            for run in paragraph.runs:
                run.font.size = Pt(20)
